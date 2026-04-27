from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "SMART_SUPPLY_CHAIN_AND_INVENTORY_MANAGEMENT_SYSTEM.pptx"
FALLBACK_OUTPUT = ROOT / "SMART_SUPPLY_CHAIN_AND_INVENTORY_MANAGEMENT_SYSTEM_UPDATED.pptx"

TITLE = "SMART SUPPLY CHAIN AND INVENTORY MANAGEMENT SYSTEM POWERED BY MICROSERVICES"
SUBTITLE = "Capstone Project Presentation"

BG = RGBColor(245, 247, 250)
NAVY = RGBColor(21, 42, 87)
BLUE = RGBColor(39, 110, 241)
TEAL = RGBColor(22, 163, 148)
ORANGE = RGBColor(245, 124, 0)
SLATE = RGBColor(74, 85, 104)
LIGHT = RGBColor(226, 232, 240)
WHITE = RGBColor(255, 255, 255)
PALE_BLUE = RGBColor(230, 238, 250)


def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_top_band(slide, slide_label, title):
    band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.9)
    )
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.color.rgb = NAVY

    label_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.18), Inches(1.3), Inches(0.3))
    lp = label_box.text_frame.paragraphs[0]
    lp.text = slide_label
    lp.font.size = Pt(12)
    lp.font.color.rgb = WHITE
    lp.font.bold = True

    title_box = slide.shapes.add_textbox(Inches(0.38), Inches(0.38), Inches(9.8), Inches(0.34))
    tp = title_box.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(23)
    tp.font.bold = True
    tp.font.color.rgb = WHITE


def add_footer(slide):
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.4), Inches(7.05), Inches(12.5), Inches(0.03)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.color.rgb = BLUE


def add_bullets(slide, items, x=0.8, y=1.4, w=11.5, h=5.2, font_size=20):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = NAVY
        p.space_after = Pt(10)


def add_box(slide, text, x, y, w, h, fill_color, font_size=16):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = fill_color
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(font_size)
    r.font.bold = True
    r.font.color.rgb = WHITE if fill_color != LIGHT else NAVY
    return shape


def add_connector(slide, x1, y1, x2, y2):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    line.line.color.rgb = SLATE
    line.line.width = Pt(2.5)


def make_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    banner = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1.0)
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = NAVY
    banner.line.color.rgb = NAVY

    badge = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(4.55), Inches(0.55), Inches(4.2), Inches(0.42)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = ORANGE
    badge.line.color.rgb = ORANGE
    bp = badge.text_frame.paragraphs[0]
    bp.alignment = PP_ALIGN.CENTER
    bp.text = "FINAL PROJECT PRESENTATION"
    bp.font.size = Pt(12)
    bp.font.bold = True
    bp.font.color.rgb = WHITE

    title = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(11.8), Inches(2.0))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = TITLE
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = NAVY

    sub = slide.shapes.add_textbox(Inches(2.0), Inches(3.6), Inches(9.3), Inches(1.0))
    p2 = sub.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    p2.text = SUBTITLE
    p2.font.size = Pt(20)
    p2.font.color.rgb = BLUE

    info = slide.shapes.add_textbox(Inches(2.0), Inches(4.5), Inches(9.3), Inches(1.4))
    p3 = info.text_frame.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    p3.text = "Built with React frontend, Node.js/Express microservices, Sequelize ORM, MySQL, JWT auth, and analytics dashboards"
    p3.font.size = Pt(16)
    p3.font.color.rgb = SLATE

    summary = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(1.15), Inches(5.3), Inches(11.0), Inches(1.05)
    )
    summary.fill.solid()
    summary.fill.fore_color.rgb = PALE_BLUE
    summary.line.color.rgb = PALE_BLUE
    sp = summary.text_frame.paragraphs[0]
    sp.alignment = PP_ALIGN.CENTER
    sp.text = "A web-based platform that helps businesses manage products, inventory, users, and customer orders efficiently with secure role-based access."
    sp.font.size = Pt(15)
    sp.font.color.rgb = NAVY


def make_content_slide(prs, slide_label, section, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide, slide_label, section)
    add_bullets(slide, bullets, y=1.35, h=5.4, font_size=19)
    add_footer(slide)


def make_architecture_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide, "Slide 3", "Overall Explanation with Design Diagram")

    add_box(slide, "Admin", 0.8, 1.6, 1.6, 0.8, ORANGE)
    add_box(slide, "Customer", 0.8, 3.0, 1.6, 0.8, TEAL)
    add_box(slide, "React Frontend\nRole-based UI\nCharts + Forms", 3.0, 2.0, 2.3, 1.5, BLUE)
    add_box(slide, "API Gateway\nPort 5000", 5.8, 2.0, 1.8, 1.2, NAVY)
    add_box(slide, "Auth Service\n5001", 8.1, 1.0, 1.7, 0.85, TEAL)
    add_box(slide, "Product Service\n5002", 8.1, 2.0, 1.7, 0.85, BLUE)
    add_box(slide, "Order Service\n5003", 8.1, 3.0, 1.7, 0.85, ORANGE)
    add_box(slide, "User Service\n5004", 8.1, 4.0, 1.7, 0.85, SLATE)
    add_box(slide, "MySQL Database\nUsers, Products, Orders,\nCart, Sessions, Settings", 10.3, 2.0, 2.2, 1.7, LIGHT, 15)

    add_connector(slide, 2.4, 2.0, 3.0, 2.4)
    add_connector(slide, 2.4, 3.4, 3.0, 3.0)
    add_connector(slide, 5.3, 2.7, 5.8, 2.6)
    add_connector(slide, 7.6, 2.6, 8.1, 1.4)
    add_connector(slide, 7.6, 2.6, 8.1, 2.4)
    add_connector(slide, 7.6, 2.6, 8.1, 3.4)
    add_connector(slide, 7.6, 2.6, 8.1, 4.4)
    add_connector(slide, 9.8, 1.4, 10.3, 2.4)
    add_connector(slide, 9.8, 2.4, 10.3, 2.6)
    add_connector(slide, 9.8, 3.4, 10.3, 2.9)
    add_connector(slide, 9.8, 4.4, 10.3, 3.2)

    bullets = [
        "Gateway routes requests to dedicated authentication, product, order, and user services.",
        "Customer workflows cover browsing, wishlist, cart, checkout, tracking, and returns.",
        "Admin workflows cover product control, approvals, users, settings, reports, and dispatch.",
        "Security layer includes JWT validation, session rotation, role checks, CSRF protection, CORS, Helmet, and rate limiting.",
    ]
    add_bullets(slide, bullets, x=0.8, y=5.15, w=11.8, h=1.8, font_size=14)
    add_footer(slide)


def make_dataset_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide, "Slide 4", "Dataset Description")

    add_box(slide, "Master Data", 0.8, 1.4, 3.7, 2.0, BLUE)
    add_box(slide, "Transactional Data", 4.8, 1.4, 3.7, 2.0, TEAL)
    add_box(slide, "System / Control Data", 8.8, 1.4, 3.7, 2.0, ORANGE)

    left = [
        "Product: name, category, price, quantity, reserved quantity, weight, image URL/key",
        "User: email, role, admin approval status, validity flag, address details",
    ]
    mid = [
        "Cart: user-product mapping with quantity",
        "Order: customer, status, total, shipping info, delivery timestamps",
        "OrderItem: product-wise quantity and price snapshot per order",
        "Wishlist: saved products for customer convenience",
    ]
    right = [
        "UserSession: refresh-token session tracking and revocation",
        "AppSetting: shipping charge, return window, cart limit, low-stock threshold, delivery ETA range",
        "AuditLog: product restock, order actions, and admin activity history",
    ]

    add_bullets(slide, left, x=0.95, y=1.95, w=3.25, h=2.0, font_size=14)
    add_bullets(slide, mid, x=4.95, y=1.95, w=3.25, h=2.2, font_size=14)
    add_bullets(slide, right, x=8.95, y=1.95, w=3.25, h=2.0, font_size=14)

    note = slide.shapes.add_textbox(Inches(0.9), Inches(4.35), Inches(11.5), Inches(1.7))
    p = note.text_frame.paragraphs[0]
    p.text = "Dataset nature: structured relational operational data stored in MySQL and modeled through Sequelize entities. The system uses live transactional records rather than a static ML dataset."
    p.font.size = Pt(18)
    p.font.color.rgb = NAVY
    add_footer(slide)


def make_tools_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide, "Slide 5", "Implementation Tools / Software")
    bullets = [
        "Frontend: React 19, React Router, Vite, Chart.js, HTML, CSS, JavaScript",
        "Backend: Node.js, Express.js, Sequelize ORM, http-proxy-middleware, Multer",
        "Database: MySQL with relational tables for users, products, orders, carts, sessions, settings, and logs",
        "Security: JWT, bcrypt, cookie-based auth, CSRF validation, Helmet, CORS, express-rate-limit",
        "Utilities: Nodemailer for emails, AWS S3 integration for product images, PM2 and Nodemon for service management",
        "Development approach: modular microservices with gateway routing on ports 5000 to 5004",
    ]
    add_bullets(slide, bullets, y=1.35, font_size=18)
    add_footer(slide)


def make_output_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide, "Slide 6", "Output")
    bullets = [
        "Secure registration, email verification, login, forgot-password, and approval workflows",
        "Customer-side product browsing with search, filtering, sorting, wishlist, cart persistence, and checkout",
        "Admin-side product CRUD, low-stock monitoring, restocking, and image management",
        "Order lifecycle handling: place, cancel, dispatch, auto-deliver, and return within configured window",
        "Dashboards and reports: revenue by month, monthly orders, best-selling product, profitable category, CSV exports, and audit logs",
        "Operational flexibility through configurable settings like shipping charge, stock threshold, and delivery timing",
    ]
    add_bullets(slide, bullets, y=1.35, font_size=18)
    add_footer(slide)


def make_references_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide, "Slide 10", "References")
    bullets = [
        "Chopra, S. and Meindl, P. Supply Chain Management: Strategy, Planning, and Operation. Pearson.",
        "Chen, J. and Paulraj, A. Towards a theory of supply chain management. Journal of Operations Management.",
        "Date, C. J. An Introduction to Database Systems. Pearson.",
        "React Documentation - https://react.dev/",
        "Express.js Documentation - https://expressjs.com/",
        "Sequelize Documentation - https://sequelize.org/ and JWT Introduction - https://jwt.io/introduction",
    ]
    add_bullets(slide, bullets, y=1.35, font_size=17)
    add_footer(slide)


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    make_title_slide(prs)
    make_content_slide(
        prs,
        "Slide 1",
        "Synopsis",
        [
            "The SMART SUPPLY CHAIN AND INVENTORY MANAGEMENT SYSTEM POWERED BY MICROSERVICES is a web-based platform designed for efficient product, stock, and order management.",
            "It combines authentication, inventory control, customer ordering, admin approvals, analytics, and configurable business rules in one centralized system.",
            "The solution is intended for small to medium-scale businesses that need improved inventory visibility, secure access control, and smoother order processing.",
        ],
    )
    make_content_slide(
        prs,
        "Slide 2",
        "Introduction",
        [
            "Small and medium-sized warehouse or retail businesses often depend on spreadsheets, manual records, or disconnected tools for stock and order handling.",
            "This causes inventory mismatch, delayed fulfillment, poor order traceability, and difficulty in making timely operational decisions.",
            "The proposed system addresses these issues through a secure full-stack application with microservice separation, real-time stock validation, and centralized role-based workflows.",
        ],
    )
    make_architecture_slide(prs)
    make_dataset_slide(prs)
    make_tools_slide(prs)
    make_output_slide(prs)
    make_content_slide(
        prs,
        "Slide 7",
        "Conclusion",
        [
            "The project successfully developed a secure and scalable inventory management and order processing system.",
            "Role-based access, protected sessions, stock reservation logic, and configurable settings improve operational safety and accuracy.",
            "The microservices-based design makes the platform modular, easier to maintain, and ready for further expansion.",
        ],
    )
    make_content_slide(
        prs,
        "Slide 8",
        "Future Enhancement",
        [
            "Integrate payment gateways for real transactions.",
            "Add cloud deployment, Docker/Kubernetes, and CI/CD automation.",
            "Connect external shipping/logistics APIs for real delivery tracking.",
            "Introduce AI-based demand forecasting and product recommendation modules.",
            "Extend notifications through email, SMS, and mobile push alerts.",
        ],
    )
    make_content_slide(
        prs,
        "Slide 9",
        "Outcome",
        [
            "A working capstone application with customer and admin modules, inventory dashboards, reports, approval flow, and secure authentication.",
            "Improved inventory accuracy through reservation tracking, low-stock visibility, order-state management, and return-based stock restoration.",
            "A strong foundation for future smart supply chain capabilities such as forecasting, logistics integration, and cloud-native deployment.",
        ],
    )
    make_references_slide(prs)

    try:
        prs.save(OUTPUT)
        return OUTPUT
    except PermissionError:
        prs.save(FALLBACK_OUTPUT)
        return FALLBACK_OUTPUT


if __name__ == "__main__":
    saved_to = build_presentation()
    print(saved_to)
